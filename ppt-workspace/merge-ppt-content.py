#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT内容合并脚本
将源PPT的内容放到模板PPT的样式框架中
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import sys
import os

# 确保输出使用UTF-8编码
sys.stdout.reconfigure(encoding='utf-8')

def extract_slide_content(slide):
    """提取幻灯片的所有文本内容"""
    content = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            content.append({
                'text': shape.text.strip(),
                'is_title': shape.name.startswith('Title') if hasattr(shape, 'name') else False,
                'left': shape.left if hasattr(shape, 'left') else 0,
                'top': shape.top if hasattr(shape, 'top') else 0,
                'width': shape.width if hasattr(shape, 'width') else 0,
                'height': shape.height if hasattr(shape, 'height') else 0,
            })
    return content

def is_placeholder_text(text):
    """判断是否为占位符文本"""
    placeholder_keywords = [
        'title', 'content', 'subtitle', '标题', '副标题', '内容',
        '点击', '替换', 'click', 'replace', 'main content',
        'template', '模板', 'chapter', '章节', 'https://star.vision'
    ]
    text_lower = text.lower()
    return any(keyword in text_lower for keyword in placeholder_keywords)

def clear_placeholder_text(slide):
    """清除幻灯片中的占位符文本"""
    cleared_count = 0
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if is_placeholder_text(shape.text):
                # 清空文本框内容
                if hasattr(shape, "text_frame"):
                    shape.text_frame.clear()
                    cleared_count += 1
    return cleared_count

def fill_content_to_slide(slide, content_list):
    """将内容填充到幻灯片中"""
    # 首先清除占位符
    clear_placeholder_text(slide)

    # 查找可用的文本框
    text_frames = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text_frames.append(shape)

    if not text_frames:
        print(f"  警告: 该页面没有可用的文本框")
        return False

    # 按位置排序内容（从上到下，从左到右）
    sorted_content = sorted(content_list, key=lambda x: (x['top'], x['left']))

    # 填充内容
    for i, content in enumerate(sorted_content):
        if i < len(text_frames):
            text_frame = text_frames[i].text_frame
            text_frame.clear()

            # 添加文本
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            p.text = content['text']

            # 如果是标题，设置较大字体
            if content['is_title'] or i == 0:
                p.font.size = Pt(32)
                p.font.bold = True
            else:
                p.font.size = Pt(18)

    return True

def merge_ppt_content(source_file, template_file, output_file):
    """合并PPT内容"""
    print(f"开始处理...")
    print(f"源文件: {source_file}")
    print(f"模板文件: {template_file}")
    print(f"输出文件: {output_file}")
    print("-" * 60)

    # 读取源PPT
    print("\n1. 读取源PPT...")
    source_prs = Presentation(source_file)
    source_slide_count = len(source_prs.slides)
    print(f"   源PPT共 {source_slide_count} 页")

    # 提取源PPT内容
    print("\n2. 提取源PPT内容...")
    source_contents = []
    for i, slide in enumerate(source_prs.slides):
        content = extract_slide_content(slide)
        source_contents.append(content)
        print(f"   第{i+1}页: 提取了 {len(content)} 个文本元素")

    # 读取模板PPT
    print("\n3. 读取模板PPT...")
    template_prs = Presentation(template_file)
    template_slide_count = len(template_prs.slides)
    print(f"   模板PPT共 {template_slide_count} 页")

    # 创建新演示文稿，复制模板
    print("\n4. 创建新演示文稿...")
    output_prs = Presentation(template_file)

    # 确保新PPT有足够的页面
    print(f"\n5. 调整页面数量...")
    current_slide_count = len(output_prs.slides)

    # 如果模板页面少于源PPT，复制最后一页布局来添加新页面
    if current_slide_count < source_slide_count:
        print(f"   需要添加 {source_slide_count - current_slide_count} 页")
        last_layout = output_prs.slides[-1].slide_layout
        for i in range(source_slide_count - current_slide_count):
            output_prs.slides.add_slide(last_layout)
    # 如果模板页面多于源PPT，删除多余页面
    elif current_slide_count > source_slide_count:
        print(f"   需要删除 {current_slide_count - source_slide_count} 页")
        # python-pptx不支持直接删除幻灯片，所以我们创建新的演示文稿
        # 只保留需要的页面
        new_prs = Presentation()
        new_prs.slide_width = output_prs.slide_width
        new_prs.slide_height = output_prs.slide_height

        # 复制slide master
        for i in range(source_slide_count):
            slide = output_prs.slides[i]
            slide_layout = slide.slide_layout
            new_slide = new_prs.slides.add_slide(slide_layout)

            # 复制所有形状
            for shape in slide.shapes:
                # 这里简化处理，保留布局
                pass

        output_prs = new_prs

    # 填充内容
    print(f"\n6. 清理占位文字并填充源内容...")
    for i, (slide, content) in enumerate(zip(output_prs.slides, source_contents)):
        print(f"   处理第{i+1}页...")

        # 清除占位符
        cleared = clear_placeholder_text(slide)
        print(f"     - 清除了 {cleared} 个占位符")

        # 填充内容
        if content:
            success = fill_content_to_slide(slide, content)
            if success:
                print(f"     - 填充了 {len(content)} 个文本元素")
        else:
            print(f"     - 该页无内容")

    # 保存
    print(f"\n7. 保存文件...")
    output_prs.save(output_file)

    # 检查文件大小
    file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
    print(f"   文件已保存: {output_file}")
    print(f"   文件大小: {file_size:.2f} MB")

    print("\n" + "=" * 60)
    print("✓ 处理完成！")
    print("=" * 60)

if __name__ == "__main__":
    # 文件路径
    source_file = "AI协同工作新范式.pptx"
    template_file = "模板1.pptx"
    output_file = "AI协同工作新范式-应用模板.pptx"

    # 执行合并
    try:
        merge_ppt_content(source_file, template_file, output_file)
    except Exception as e:
        print(f"\n❌ 错误: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
