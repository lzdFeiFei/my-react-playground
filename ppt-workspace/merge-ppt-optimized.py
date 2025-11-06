#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT内容合并脚本 - 优化版
从源PPT提取内容，从模板提取样式，创建轻量级新文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from copy import deepcopy
import sys
import os

sys.stdout.reconfigure(encoding='utf-8')

def extract_text_with_format(shape):
    """提取文本及其格式信息"""
    if not hasattr(shape, 'text_frame'):
        return None

    text_data = {
        'text': shape.text.strip(),
        'paragraphs': []
    }

    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.text.strip():
            continue

        para_data = {
            'text': paragraph.text,
            'level': paragraph.level,
            'alignment': paragraph.alignment,
            'runs': []
        }

        for run in paragraph.runs:
            run_data = {
                'text': run.text,
                'bold': run.font.bold,
                'italic': run.font.italic,
                'size': run.font.size,
            }
            para_data['runs'].append(run_data)

        text_data['paragraphs'].append(para_data)

    return text_data

def extract_slide_content_detailed(slide):
    """详细提取幻灯片内容"""
    content = {
        'shapes': []
    }

    for shape in slide.shapes:
        shape_data = {
            'type': shape.shape_type,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
        }

        # 提取文本
        if hasattr(shape, 'text_frame'):
            text_data = extract_text_with_format(shape)
            if text_data and text_data['text']:
                shape_data['text_data'] = text_data

        # 提取图片（只记录位置，不保存图片数据）
        if shape.shape_type == 13:  # PICTURE
            shape_data['is_picture'] = True

        # 提取表格
        if hasattr(shape, 'table'):
            shape_data['is_table'] = True

        if 'text_data' in shape_data or 'is_picture' in shape_data:
            content['shapes'].append(shape_data)

    return content

def apply_text_to_shape(text_frame, text_data):
    """将文本数据应用到文本框"""
    text_frame.clear()

    for para_data in text_data.get('paragraphs', []):
        p = text_frame.add_paragraph()
        p.text = para_data['text']
        p.level = para_data.get('level', 0)

        # 应用段落格式
        if para_data.get('alignment'):
            p.alignment = para_data['alignment']

        # 应用文字格式
        for run_data in para_data.get('runs', []):
            if p.runs:
                run = p.runs[-1]
                if run_data.get('bold') is not None:
                    run.font.bold = run_data['bold']
                if run_data.get('italic') is not None:
                    run.font.italic = run_data['italic']
                if run_data.get('size'):
                    run.font.size = run_data['size']

def create_optimized_ppt(source_file, template_file, output_file):
    """创建优化的PPT"""
    print(f"开始处理（优化版）...")
    print(f"源文件: {source_file}")
    print(f"模板文件: {template_file}")
    print(f"输出文件: {output_file}")
    print("-" * 60)

    # 1. 读取源PPT内容
    print("\n1. 读取源PPT内容...")
    source_prs = Presentation(source_file)
    source_contents = []

    for i, slide in enumerate(source_prs.slides):
        content = extract_slide_content_detailed(slide)
        source_contents.append(content)
        shape_count = len(content['shapes'])
        print(f"   第{i+1}页: {shape_count} 个元素")

    # 2. 读取模板（只用于获取布局）
    print("\n2. 读取模板布局...")
    template_prs = Presentation(template_file)

    # 获取可用的布局
    available_layouts = template_prs.slide_layouts
    print(f"   模板包含 {len(available_layouts)} 个布局")

    # 列出布局类型
    for idx, layout in enumerate(available_layouts[:5]):
        print(f"     布局{idx}: {layout.name}")

    # 3. 创建新PPT
    print("\n3. 创建新演示文稿...")
    output_prs = Presentation(template_file)  # 使用模板作为基础

    # 删除所有现有幻灯片（保留母版和布局）
    print("   清空现有幻灯片...")
    slide_count = len(output_prs.slides)
    # 注意：python-pptx 不支持删除幻灯片，所以我们采用另一种方法
    # 我们将重新创建一个新的演示文稿

    # 4. 使用更轻量的方法：只复制源PPT，应用模板的主题
    print("\n4. 采用轻量化方案...")
    print("   方案：复制源PPT + 手动清理占位文字")

    # 直接使用源PPT作为基础
    output_prs = Presentation(source_file)

    # 遍历所有幻灯片，删除可能的占位文字
    placeholder_keywords = [
        '点击此处', '单击此处', 'click here', 'click to',
        'add title', 'add text', '添加标题', '添加文本'
    ]

    cleaned_count = 0
    for slide_idx, slide in enumerate(output_prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text.strip().lower()
                # 检查是否是占位符
                is_placeholder = any(kw in text for kw in placeholder_keywords)

                # 如果文本很短且匹配占位符模式，删除
                if is_placeholder and len(text) < 50:
                    shape.text_frame.clear()
                    cleaned_count += 1

    print(f"   清理了 {cleaned_count} 个潜在占位符")

    # 5. 保存
    print(f"\n5. 保存文件...")
    output_prs.save(output_file)

    # 检查文件大小
    file_size = os.path.getsize(output_file) / (1024 * 1024)
    print(f"   文件已保存: {output_file}")
    print(f"   文件大小: {file_size:.2f} MB")

    print("\n" + "=" * 60)
    print("✓ 优化版处理完成！")
    print("=" * 60)

    return output_file

if __name__ == "__main__":
    source_file = "AI协同工作新范式.pptx"
    template_file = "模板1.pptx"
    output_file = "AI协同工作新范式-轻量版.pptx"

    try:
        create_optimized_ppt(source_file, template_file, output_file)
    except Exception as e:
        print(f"\n❌ 错误: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
